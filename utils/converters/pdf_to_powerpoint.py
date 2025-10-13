import asyncio
import os
import tempfile
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from supabase import create_client

supabase_url = os.getenv("SUPABASE_URL")
supabase_key = os.getenv("SUPABASE_KEY")
supabase = create_client(supabase_url, supabase_key)

BUCKET_NAME = "docsbucket"


async def convert_pdf_to_powerpoint(file):
    contents = await file.read()
    def _convert():
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(contents)
            tmp_path = tmp.name
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Extract text from PDF using pdfplumber
        with pdfplumber.open(tmp_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                # Extract text
                text = page.extract_text()
                
                # Extract tables if any
                tables = page.extract_tables()
                
                # Add blank slide
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title with page number
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
                title_frame = title_box.text_frame
                title_frame.text = f"Page {page_num}"
                title_frame.paragraphs[0].font.size = Pt(28)
                title_frame.paragraphs[0].font.bold = True
                
                # Add main text content
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5.5)
                
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                
                if text:
                    text_frame.text = text.strip()
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                
                # Add table if exists
                if tables:
                    # Create new slide for table
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    # Add table title
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
                    title_frame = title_box.text_frame
                    title_frame.text = f"Page {page_num} - Table"
                    title_frame.paragraphs[0].font.size = Pt(28)
                    title_frame.paragraphs[0].font.bold = True
                    
                    # Add first table only (to avoid overcrowding)
                    table_data = tables[0]
                    rows = len(table_data)
                    cols = len(table_data[0]) if rows > 0 else 0
                    
                    if rows > 0 and cols > 0:
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(9)
                        height = Inches(5.5)
                        
                        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
                        table = shape.table
                        
                        # Add table data
                        for row_idx, row in enumerate(table_data):
                            for col_idx, cell_value in enumerate(row):
                                cell = table.cell(row_idx, col_idx)
                                cell.text = str(cell_value) if cell_value else ""
                                
                                # Format header row
                                if row_idx == 0:
                                    cell.text_frame.paragraphs[0].font.bold = True
        
        pptx_path = tmp_path.replace(".pdf", ".pptx")
        prs.save(pptx_path)
        
        return tmp_path, pptx_path
    
    tmp_path, pptx_path = await asyncio.to_thread(_convert)
    
    try:
        file_name = os.path.basename(pptx_path)
        with open(pptx_path, "rb") as pptx_file:
            supabase.storage.from_(BUCKET_NAME).upload(
                file_name,
                pptx_file,
                {"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "upsert": "false"}
            )
        download_url = f"{supabase_url}/storage/v1/object/public/{BUCKET_NAME}/{file_name}"
    except Exception as e:
        print(f"❌ Supabase upload failed: {str(e)}")
        raise
    finally:
        os.remove(tmp_path)
        os.remove(pptx_path)
    
    return download_url, file_name